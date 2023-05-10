import sys
from pptx import Presentation
from gtts import gTTS
import os
def text_to_speech(text, output_file):
    tts = gTTS(text=text, lang='ro')
    tts.save(output_file)

def extract_text_from_ppt(file_path):
    presentation = Presentation(file_path)
    text = ""

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + " "

    return text

if __name__ == "__main__":
    if len(sys.argv) > 1:
        ppt_file = sys.argv[1]
        extracted_text = extract_text_from_ppt(ppt_file)
        
        output_file = "output.mp3"
        text_to_speech(extracted_text, output_file)
        
        os.system(f"start {output_file}")
    else:
        print("Please provide a PowerPoint file path.")